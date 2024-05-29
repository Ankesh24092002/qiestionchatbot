from flask import Flask, render_template, request, session, jsonify, current_app
from flask_session import Session
import os
import openai
import cassio
from langchain_astradb import AstraDBVectorStore
from langchain_community.llms import OpenAI
from langchain.text_splitter import CharacterTextSplitter
from langchain_openai import OpenAIEmbeddings
from langchain_openai import AzureOpenAIEmbeddings
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
from dotenv import load_dotenv
from urllib.request import urlopen
from bs4 import BeautifulSoup
import requests
from astrapy.db import AstraDB
import secrets
import random
import string
import datetime
import uuid
from azure.cosmos import CosmosClient, exceptions, PartitionKey

load_dotenv(override=True)
COSMOS_DB_URI = os.getenv('COSMOS_DB_URI')
COSMOS_DB_KEY = os.getenv('COSMOS_DB_KEY')
COSMOS_DB_NAME = os.getenv('COSMOS_DB_NAME')
COSMOS_DB_CONTAINER = os.getenv('COSMOS_DB_CONTAINER')

client = CosmosClient(COSMOS_DB_URI, COSMOS_DB_KEY)
database = client.create_database_if_not_exists(id=COSMOS_DB_NAME)
container = database.create_container_if_not_exists(
    id=COSMOS_DB_CONTAINER,
    partition_key=PartitionKey(path="/user_id"),
    offer_throughput=400
)

app = Flask(__name__)
app.config["SESSION_PERMANENT"] = False
app.config["SESSION_TYPE"] = "filesystem"
app.config["AstraVectorStore"] = None
Session(app)

load_dotenv(override=True)
from openai import AzureOpenAI

client = AzureOpenAI(
    azure_endpoint=os.getenv("AZURE_OPENAI_ENDPOINT"),
    api_key=os.getenv("AZURE_OPENAI_KEY"),
    api_version="2024-02-15-preview"
)

ASTRA_DB_APPLICATION_TOKEN = os.getenv('ASTRA_DB_APPLICATION_TOKEN')
ASTR_DB_ID = os.getenv('ASTR_DB_ID')
ASTRA_DB_API_ENDPOINT = os.getenv('ASTRA_DB_API_ENDPOINT')
AZURE_OPENAI_KEY = os.getenv('AZURE_OPENAI_KEY')
AZURE_OPENAI_ENDPOINT = os.getenv('AZURE_OPENAI_ENDPOINT')
app.secret_key = os.getenv('SECRET_KEY', 'supersecretkey')

def generate_random_string(length=10):
    characters = string.ascii_lowercase + string.digits + '_'
    random_string = ''.join(random.choice(characters) for _ in range(length - 2))
    return random.choice(string.ascii_lowercase) + random_string + random.choice(string.ascii_lowercase)

def initialize_astra_vector_store(table_name):
    cassio.init(token=ASTRA_DB_APPLICATION_TOKEN, database_id=ASTR_DB_ID)
    embedding = AzureOpenAIEmbeddings(azure_deployment="text-embedding-ada-002", api_key=AZURE_OPENAI_KEY, azure_endpoint=AZURE_OPENAI_ENDPOINT)
    astra_vector_store = AstraDBVectorStore(
        embedding=embedding,
        collection_name=table_name,
        api_endpoint=ASTRA_DB_API_ENDPOINT,
        token=ASTRA_DB_APPLICATION_TOKEN
    )
    return astra_vector_store

def delete_collection(table_name):
    db = AstraDB(
        token=ASTRA_DB_APPLICATION_TOKEN,
        api_endpoint=ASTRA_DB_API_ENDPOINT,
    )
    db.delete_collection(collection_name=table_name)

def preprocessor_files(uploaded_file):
    texts = []
    if uploaded_file:
        if uploaded_file.filename.endswith('.pdf'):
            texts.extend(preprocess_pdf(uploaded_file))
        elif uploaded_file.filename.endswith(('.doc', '.docx')):
            texts.extend(preprocess_word(uploaded_file))
        elif uploaded_file.filename.endswith('.txt'):
            texts.extend(preprocess_text(uploaded_file))
        elif uploaded_file.filename.endswith('.md'):
            texts.extend(preprocess_markdown(uploaded_file))
        elif uploaded_file.filename.endswith(('.html', '.htm')):
            texts.extend(preprocess_html(uploaded_file))
        elif uploaded_file.filename.endswith('.pptx'):
            texts.extend(preprocess_pptx(uploaded_file))
    return texts

def preprocess_pptx(uploaded_file):
    prs = Presentation(uploaded_file)
    raw_text = ''
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                raw_text += shape.text + '\n'
    text_splitter = CharacterTextSplitter(
        separator='\n',
        chunk_size=900,
        chunk_overlap=200,
        length_function=len
    )
    texts = text_splitter.split_text(raw_text)
    return texts

def preprocess_text(uploaded_file):
    raw_text = uploaded_file.read().decode('utf-8')
    text_splitter = CharacterTextSplitter(
        separator='\n',
        chunk_size=900,
        chunk_overlap=200,
        length_function=len
    )
    texts = text_splitter.split_text(raw_text)
    return texts

def preprocess_markdown(uploaded_file):
    raw_text = uploaded_file.read().decode('utf-8')
    text_splitter = CharacterTextSplitter(
        separator='\n',
        chunk_size=900,
        chunk_overlap=200,
        length_function=len
    )
    texts = text_splitter.split_text(raw_text)
    return texts

def preprocess_html(uploaded_file):
    raw_text = uploaded_file.read().decode('utf-8')
    text_splitter = CharacterTextSplitter(
        separator='\n',
        chunk_size=900,
        chunk_overlap=200,
        length_function=len
    )
    texts = text_splitter.split_text(raw_text)
    return texts

def preprocess_pdf(uploaded_file):
    pdf_reader = PdfReader(uploaded_file)
    raw_text = ''
    for page in pdf_reader.pages:
        content = page.extract_text()
        if content:
            raw_text += content
    text_splitter = CharacterTextSplitter(
        separator='\n',
        chunk_size=900,
        chunk_overlap=200,
        length_function=len
    )
    texts = text_splitter.split_text(raw_text)
    return texts

def preprocess_word(uploaded_file):
    document = Document(uploaded_file)
    raw_text = ''
    for paragraph in document.paragraphs:
        raw_text += paragraph.text
    text_splitter = CharacterTextSplitter(
        separator='\n',
        chunk_size=900,
        chunk_overlap=200,
        length_function=len
    )
    texts = text_splitter.split_text(raw_text)
    return texts

def perform_query(query_text, astra_vector_store):
    vectorDB_answer = astra_vector_store.similarity_search_with_score(query_text, k=1)
    res, score = vectorDB_answer[0]
    return res.page_content, score

def perform_query_chat(message_history):
    response = client.chat.completions.create(
        model="gpt35turbo16k",
        messages=message_history,
        temperature=0.7,
        max_tokens=800,
        top_p=0.95,
        frequency_penalty=0,
        presence_penalty=0,
        stop=None
    )
    return response

@app.route('/')
def index():
    return render_template('newfinal.html', uploaded_files=session.get('uploaded_files', []))

message_history = []

@app.route('/upload', methods=['POST'])
def upload():
    if 'upload_file' in request.files:
        file = request.files['upload_file']
        if file.filename != '':
            if session.get('table_name'):
                table_name = session.get('table_name')
                delete_collection(table_name)
                session['table_name'] = None

            table_name = generate_random_string()
            session['table_name'] = table_name
            session['uploaded_files'] = [file.filename]
            texts = preprocessor_files(file)
            astra_vector_store = initialize_astra_vector_store(table_name)
            current_app.config["AstraVectorStore"] = astra_vector_store
            astra_vector_store.add_texts(texts)

            # Set a user_id in the session if it doesn't exist
            if 'user_id' not in session:
                session['user_id'] = generate_random_string()

            return "File uploaded and preprocessed successfully!"
    return "No file selected!"

@app.route('/chat', methods=['POST'])
def chatbot():
    user_message = request.form.get('user_message')

    if user_message:
        if session.get('uploaded_files') and current_app.config["AstraVectorStore"] is not None:
            astra_vector_store = current_app.config["AstraVectorStore"]
            vectorDB_answer, score = perform_query(user_message, astra_vector_store)

            prompt = f"""
                You are an AI assistant that helps people generate relevant content based on the context provided.
                
                Context:
                {vectorDB_answer}
                
                User's query:
                {user_message}
                
                Please generate a response based on the above context and query.
            """

            message_history.append({"role": "user", "content": prompt})
            response = perform_query_chat(message_history)

            message_history.append({"role": "assistant", "content": response.choices[0].message.content})

            # Save the response to Cosmos DB
            try:
                container.create_item(body={
                    "id": str(uuid.uuid4()),  # Unique identifier for the document
                    "user_id": session['user_id'],
                    "user_message": user_message,
                    "response": response.choices[0].message.content,
                    "timestamp": datetime.datetime.utcnow().isoformat()
                })
            except exceptions.CosmosHttpResponseError as e:
                return jsonify({"response": "Failed to store the response in Cosmos DB!"})

            return jsonify({"response": response.choices[0].message.content})
        else:
            return jsonify({"response": "Please upload a file to start the conversation!"})
    return jsonify({"response": "Please provide a message!"})

@app.route('/delete', methods=['POST'])
def delete():
    table_name = session.get('table_name')
    if table_name:
        delete_collection(table_name)
        session.pop('table_name', None)
        session.pop('uploaded_files', None)
        return jsonify({"response": "All uploaded files and vectors have been deleted!"})
    return jsonify({"response": "No files to delete!"})

if __name__ == '__main__':
    app.run(debug=True)
